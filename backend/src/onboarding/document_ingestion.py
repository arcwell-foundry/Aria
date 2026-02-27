"""Document Upload & Ingestion Pipeline for onboarding (US-904).

Handles document upload, parsing, chunking, and knowledge extraction.
Pipeline:
    1. Upload to Supabase Storage
    2. Format detection & text extraction
    3. Semantic chunking (structure-aware)
    4. Entity extraction via LLM
    5. Embedding generation (pgvector)
    6. Source quality scoring
    7. Knowledge extraction → Corporate Memory
"""

import asyncio
import json
import logging
import uuid
from datetime import UTC, datetime
from io import BytesIO
from typing import Any

from src.core.llm import LLMClient
from src.core.task_types import TaskType
from src.db.supabase import SupabaseClient

logger = logging.getLogger(__name__)

# File size limits
MAX_FILE_SIZE = 50 * 1024 * 1024  # 50MB per file
MAX_COMPANY_TOTAL = 500 * 1024 * 1024  # 500MB per company

SUPPORTED_TYPES: dict[str, str] = {
    "application/pdf": "pdf",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": "docx",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": "pptx",
    "text/plain": "txt",
    "text/markdown": "md",
    "text/csv": "csv",
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": "xlsx",
    "image/png": "image",
    "image/jpeg": "image",
    "image/webp": "image",
}


class DocumentIngestionService:
    """Handles document upload, parsing, chunking, and knowledge extraction.

    Pipeline:
        1. Upload to Supabase Storage
        2. Format detection & text extraction
        3. Semantic chunking (structure-aware)
        4. Entity extraction via LLM
        5. Embedding generation (pgvector)
        6. Source quality scoring
        7. Knowledge extraction → Corporate Memory
    """

    def __init__(self) -> None:
        """Initialize with Supabase client and LLM client."""
        self._db = SupabaseClient.get_client()
        self._llm = LLMClient()

    async def validate_upload(
        self,
        company_id: str,
        filename: str,  # noqa: ARG002
        file_size: int,
        content_type: str,
    ) -> dict[str, Any]:
        """Validate a file before upload.

        Args:
            company_id: The company's UUID.
            filename: Original filename.
            file_size: File size in bytes.
            content_type: MIME content type.

        Returns:
            Dict with keys: valid (bool), reason (str | None), file_type (str).
        """
        file_type = SUPPORTED_TYPES.get(content_type)
        if not file_type:
            return {
                "valid": False,
                "reason": f"Unsupported file type: {content_type}",
                "file_type": "",
            }

        if file_size > MAX_FILE_SIZE:
            return {
                "valid": False,
                "reason": "File too large. Maximum is 50MB.",
                "file_type": file_type,
            }

        # Check company total storage usage
        result = (
            self._db.table("company_documents")
            .select("file_size_bytes")
            .eq("company_id", company_id)
            .execute()
        )
        current_total = sum(doc["file_size_bytes"] for doc in (result.data or []))
        if current_total + file_size > MAX_COMPANY_TOTAL:
            return {
                "valid": False,
                "reason": "Company storage limit reached (500MB).",
                "file_type": file_type,
            }

        return {"valid": True, "reason": None, "file_type": file_type}

    async def upload_and_process(
        self,
        company_id: str,
        user_id: str,
        filename: str,
        file_content: bytes,
        content_type: str,
    ) -> dict[str, Any]:
        """Upload file to storage and start processing pipeline.

        Args:
            company_id: The company's UUID.
            user_id: The uploading user's UUID.
            filename: Original filename.
            file_content: Raw file bytes.
            content_type: MIME content type.

        Returns:
            Document record with processing status.
        """
        file_type = SUPPORTED_TYPES.get(content_type, "unknown")

        # 1. Upload to Supabase Storage
        storage_path = f"companies/{company_id}/documents/{filename}"
        self._db.storage.from_("documents").upload(
            storage_path, file_content, {"content-type": content_type}
        )

        # 2. Create document record
        doc_record = {
            "company_id": company_id,
            "uploaded_by": user_id,
            "filename": filename,
            "file_type": file_type,
            "file_size_bytes": len(file_content),
            "storage_path": storage_path,
            "processing_status": "processing",
        }
        result = self._db.table("company_documents").insert(doc_record).execute()
        doc: dict[str, Any] = result.data[0]

        # 3. Process asynchronously (fire-and-forget)
        asyncio.create_task(
            self._process_document(doc["id"], file_content, file_type, company_id, user_id)
        )

        return doc

    async def get_company_documents(self, company_id: str) -> list[dict[str, Any]]:
        """Get all documents for a company.

        Args:
            company_id: The company's UUID.

        Returns:
            List of document records ordered by creation date (newest first).
        """
        result = (
            self._db.table("company_documents")
            .select("*")
            .eq("company_id", company_id)
            .order("created_at", desc=True)
            .execute()
        )
        return result.data or []

    # --- Internal pipeline stages ---

    async def _process_document(
        self,
        doc_id: str,
        content: bytes,
        file_type: str,
        company_id: str,
        user_id: str,
    ) -> None:
        """Run the full processing pipeline for a document.

        Args:
            doc_id: Document UUID.
            content: Raw file bytes.
            file_type: Detected file type string.
            company_id: Company UUID.
            user_id: Uploading user UUID.
        """
        try:
            # Step 1: Extract text
            await self._update_progress(doc_id, 10, "processing")
            text = await self._extract_text(content, file_type)

            if not text or len(text.strip()) < 50:
                await self._update_progress(doc_id, 100, "complete")
                logger.warning(
                    "Document had insufficient text content",
                    extra={"document_id": doc_id},
                )
                return

            # Step 2: Semantic chunking
            await self._update_progress(doc_id, 30, "processing")
            chunks = self._semantic_chunk(text)

            # Step 3: Entity extraction + embedding for each chunk
            await self._update_progress(doc_id, 50, "processing")
            total_entities: list[dict[str, str]] = []
            for i, chunk in enumerate(chunks):
                entities = await self._extract_entities(chunk["content"])
                embedding = await self._generate_embedding(chunk["content"])

                self._db.table("document_chunks").insert(
                    {
                        "document_id": doc_id,
                        "chunk_index": i,
                        "content": chunk["content"],
                        "chunk_type": chunk["type"],
                        "embedding": embedding,
                        "entities": entities,
                        "metadata": chunk.get("metadata", {}),
                    }
                ).execute()
                total_entities.extend(entities)

            # Step 4: Quality scoring
            await self._update_progress(doc_id, 80, "processing")
            quality_score = await self._score_quality(text, file_type)

            # Step 4b: Seed extracted entities to Graphiti knowledge graph
            await self._seed_entities_to_graph(total_entities, company_id, user_id, doc_id)

            # Step 5: Knowledge extraction → Corporate Memory
            await self._update_progress(doc_id, 90, "processing")
            await self._extract_knowledge(text, company_id, user_id, doc_id)

            # Step 6: Finalize
            unique_entity_count = len({e.get("name", "") for e in total_entities})
            self._db.table("company_documents").update(
                {
                    "processing_status": "complete",
                    "processing_progress": 100,
                    "chunk_count": len(chunks),
                    "entity_count": unique_entity_count,
                    "quality_score": quality_score,
                }
            ).eq("id", doc_id).execute()

            logger.info(
                "Document processing complete",
                extra={
                    "document_id": doc_id,
                    "chunks": len(chunks),
                    "entities": unique_entity_count,
                    "quality_score": quality_score,
                },
            )

            # Update readiness score
            await self._update_readiness(user_id, quality_score)

            # Record episodic memory
            await self._record_episodic(
                user_id, doc_id, len(chunks), unique_entity_count, quality_score
            )

            # Record activity for feed
            try:
                from src.services.activity_service import ActivityService

                # Fetch filename from document record
                doc_result = (
                    self._db.table("company_documents")
                    .select("filename")
                    .eq("id", doc_id)
                    .maybe_single()
                    .execute()
                )
                filename = (doc_result.data or {}).get("filename", "document")

                await ActivityService().record(
                    user_id=user_id,
                    agent="analyst",
                    activity_type="document_processed",
                    title=f"Processed {filename}",
                    description=(f"ARIA processed {filename} — extracted key entities and facts"),
                    confidence=0.8,
                    related_entity_type="document",
                    related_entity_id=doc_id,
                )
            except Exception as e:
                logger.warning("Failed to record document activity: %s", e)

        except Exception:
            logger.exception(
                "Document processing failed",
                extra={"document_id": doc_id},
            )
            self._db.table("company_documents").update({"processing_status": "failed"}).eq(
                "id", doc_id
            ).execute()

    async def _extract_text(self, content: bytes, file_type: str) -> str:
        """Extract text from document based on file type.

        Args:
            content: Raw file bytes.
            file_type: Detected file type string.

        Returns:
            Extracted plain text.
        """
        if file_type in ("txt", "md", "csv"):
            return content.decode("utf-8", errors="replace")

        if file_type == "pdf":
            return self._extract_pdf(content)

        if file_type == "docx":
            return self._extract_docx(content)

        if file_type == "pptx":
            return self._extract_pptx(content)

        if file_type == "xlsx":
            return self._extract_xlsx(content)

        if file_type == "image":
            return self._extract_image_ocr(content)

        return ""

    @staticmethod
    def _extract_pdf(content: bytes) -> str:
        """Extract text from a PDF using PyMuPDF, with OCR fallback for scanned pages.

        For each page, tries native text extraction first. If a page yields
        minimal text (< 50 chars), falls back to OCR via pytesseract on a
        rendered pixmap of the page.

        Args:
            content: Raw PDF bytes.

        Returns:
            Extracted text, or empty string if parser unavailable.
        """
        try:
            import fitz  # PyMuPDF
        except ImportError:
            logger.warning("PyMuPDF not installed — PDF text extraction unavailable")
            return ""

        doc = fitz.open(stream=content, filetype="pdf")
        text = ""
        for page in doc:
            page_text = page.get_text()
            if len(page_text.strip()) < 50:
                # Scanned/image-based page — try OCR
                page_text = DocumentIngestionService._ocr_pdf_page(page)
            text += page_text + "\n\n"
        doc.close()
        return text

    @staticmethod
    def _ocr_pdf_page(page: "Any") -> str:
        """Run OCR on a single PDF page rendered as an image.

        Args:
            page: A PyMuPDF page object.

        Returns:
            OCR'd text, or empty string if Tesseract is unavailable.
        """
        try:
            import pytesseract
            from PIL import Image

            pix = page.get_pixmap(dpi=300)
            img = Image.frombytes("RGB", (pix.width, pix.height), pix.samples)
            return pytesseract.image_to_string(img)
        except ImportError:
            logger.warning(
                "pytesseract/Pillow not installed — OCR unavailable for scanned PDF page"
            )
            return ""
        except Exception as e:
            logger.warning("OCR failed for PDF page: %s", e)
            return ""

    @staticmethod
    def _extract_image_ocr(content: bytes) -> str:
        """Extract text from a standalone image file via OCR.

        Args:
            content: Raw image bytes (PNG, JPG, WebP).

        Returns:
            OCR'd text, or empty string if Tesseract is unavailable.
        """
        try:
            import pytesseract
            from PIL import Image

            img = Image.open(BytesIO(content))
            return pytesseract.image_to_string(img)
        except ImportError:
            logger.warning("pytesseract/Pillow not installed — image OCR unavailable")
            return ""
        except Exception as e:
            logger.warning("Image OCR failed: %s", e)
            return ""

    @staticmethod
    def _extract_docx(content: bytes) -> str:
        """Extract text from a DOCX file.

        Args:
            content: Raw DOCX bytes.

        Returns:
            Extracted text, or empty string if parser unavailable.
        """
        try:
            import docx

            doc = docx.Document(BytesIO(content))
            return "\n\n".join(p.text for p in doc.paragraphs if p.text.strip())
        except ImportError:
            logger.warning("python-docx not installed — DOCX text extraction unavailable")
            return ""

    @staticmethod
    def _extract_pptx(content: bytes) -> str:
        """Extract text from a PPTX file.

        Args:
            content: Raw PPTX bytes.

        Returns:
            Extracted text, or empty string if parser unavailable.
        """
        try:
            from pptx import Presentation

            prs = Presentation(BytesIO(content))
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
                text += "\n\n"
            return text
        except ImportError:
            logger.warning("python-pptx not installed — PPTX text extraction unavailable")
            return ""

    @staticmethod
    def _extract_xlsx(content: bytes) -> str:
        """Extract text from an XLSX file.

        Args:
            content: Raw XLSX bytes.

        Returns:
            Extracted text, or empty string if parser unavailable.
        """
        try:
            import openpyxl

            wb = openpyxl.load_workbook(BytesIO(content), read_only=True)
            text = ""
            for sheet in wb.worksheets:
                for row in sheet.iter_rows(values_only=True):
                    row_text = " | ".join(str(cell) for cell in row if cell)
                    if row_text:
                        text += row_text + "\n"
                text += "\n\n"
            wb.close()
            return text
        except ImportError:
            logger.warning("openpyxl not installed — XLSX text extraction unavailable")
            return ""

    @staticmethod
    def _semantic_chunk(text: str, max_chunk_size: int = 1500) -> list[dict[str, Any]]:
        """Split text into semantic chunks respecting document structure.

        Not naive character splitting — respects headers, paragraphs, tables.

        Args:
            text: Full document text.
            max_chunk_size: Maximum characters per chunk.

        Returns:
            List of dicts with 'content', 'type', and optional 'metadata'.
        """
        chunks: list[dict[str, Any]] = []
        paragraphs = [p.strip() for p in text.split("\n\n") if p.strip()]

        current_chunk = ""
        current_type = "paragraph"

        for para in paragraphs:
            # Detect chunk type
            if para.startswith("#") or (len(para) < 100 and para.isupper()):
                chunk_type = "header"
            elif "|" in para and para.count("|") > 2:
                chunk_type = "table"
            elif para.startswith(("- ", "* ", "1.", "\u2022")):
                chunk_type = "list"
            else:
                chunk_type = "paragraph"

            # If adding this would exceed limit, flush current
            if len(current_chunk) + len(para) > max_chunk_size and current_chunk:
                chunks.append(
                    {"content": current_chunk.strip(), "type": current_type, "metadata": {}}
                )
                current_chunk = ""
                current_type = chunk_type

            current_chunk += para + "\n\n"
            if chunk_type != "paragraph":
                current_type = chunk_type

        # Flush remaining
        if current_chunk.strip():
            chunks.append({"content": current_chunk.strip(), "type": current_type, "metadata": {}})

        return chunks

    async def _extract_entities(self, text: str) -> list[dict[str, str]]:
        """Extract named entities from a chunk using LLM.

        Args:
            text: Text chunk to analyze.

        Returns:
            List of entity dicts with 'name' and 'type' keys.
        """
        if len(text) < 30:
            return []

        prompt = (
            "Extract named entities from this text. Return a JSON array.\n\n"
            f"Text: {text[:2000]}\n\n"
            'Format: [{"name": "entity name", "type": "company|person|product|'
            'therapeutic_area|modality|technology|location"}]\n\n'
            "Only include clearly identifiable entities. Be precise."
        )

        try:
            response = await self._llm.generate_response(
                messages=[{"role": "user", "content": prompt}],
                max_tokens=500,
                temperature=0.2,
                task=TaskType.ENTITY_EXTRACT,
            )
            return json.loads(response)
        except (json.JSONDecodeError, Exception) as e:
            logger.debug("Entity extraction parse error: %s", e)
            return []

    async def _generate_embedding(self, text: str) -> list[float]:
        """Generate embedding vector for a text chunk using OpenAI.

        Uses the text-embedding-3-small model (same as Graphiti) to produce
        1536-dimensional vectors stored in pgvector for semantic search.

        Args:
            text: Text to embed.

        Returns:
            1536-dimensional float vector. Returns zeros if the embedding
            API is not configured or the call fails.
        """
        try:
            from openai import AsyncOpenAI

            from src.core.config import get_settings

            api_key = get_settings().OPENAI_API_KEY.get_secret_value()
            if not api_key:
                logger.debug("OpenAI API key not configured, returning zero vector")
                return [0.0] * 1536

            client = AsyncOpenAI(api_key=api_key)
            response = await client.embeddings.create(
                model="text-embedding-3-small",
                input=text[:8000],  # Limit input to stay within token bounds
            )
            return response.data[0].embedding
        except Exception as e:
            logger.warning("Embedding generation failed, returning zero vector: %s", e)
            return [0.0] * 1536

    async def _score_quality(self, text: str, file_type: str) -> float:
        """Score document source quality.

        Capabilities deck > product sheet > org chart > generic industry report.

        Args:
            text: Full document text.
            file_type: Detected file type.

        Returns:
            Quality score 0-100.
        """
        prompt = (
            "Rate the intelligence value of this document for a sales team "
            "on a 0-100 scale.\n\n"
            f"Document type: {file_type}\n"
            f"Content preview: {text[:1000]}\n\n"
            "Scoring guide:\n"
            "- 90-100: Capabilities deck, competitive analysis, pricing sheet, "
            "product roadmap\n"
            "- 70-89: Product documentation, org chart, case studies, partnership "
            "announcements\n"
            "- 50-69: General company overview, press releases, marketing materials\n"
            "- 30-49: Industry reports, generic content\n"
            "- 0-29: Irrelevant or duplicate content\n\n"
            "Respond with just the number."
        )

        try:
            response = await self._llm.generate_response(
                messages=[{"role": "user", "content": prompt}],
                max_tokens=10,
                temperature=0.1,
                task=TaskType.ENTITY_EXTRACT,
            )
            return float(response.strip())
        except (ValueError, Exception):
            return 50.0

    async def _extract_knowledge(
        self, text: str, company_id: str, user_id: str, doc_id: str
    ) -> None:
        """Extract structured knowledge from document and store in memory.

        Args:
            text: Full document text.
            company_id: Company UUID.
            user_id: User UUID.
            doc_id: Document UUID.
        """
        prompt = (
            "Extract key business facts from this document.\n\n"
            f"{text[:4000]}\n\n"
            "Return JSON array of facts:\n"
            "[\n"
            '  {"fact": "statement", "category": "product|pipeline|leadership|'
            'financial|partnership|regulatory|competitive|manufacturing", '
            '"confidence": 0.7}\n'
            "]\n\n"
            "Extract 5-15 most important facts. Be specific with names, "
            "dates, numbers."
        )

        try:
            response = await self._llm.generate_response(
                messages=[{"role": "user", "content": prompt}],
                max_tokens=2000,
                temperature=0.3,
                task=TaskType.ENTITY_EXTRACT,
            )
            facts: list[dict[str, Any]] = json.loads(response)
            for fact in facts:
                self._db.table("memory_semantic").insert(
                    {
                        "user_id": user_id,
                        "fact": fact["fact"],
                        "confidence": fact.get("confidence", 0.75),
                        "source": "document_upload",
                        "metadata": {
                            "category": fact.get("category", "general"),
                            "document_id": doc_id,
                            "company_id": company_id,
                        },
                    }
                ).execute()

            # Generate causal hypotheses from extracted facts (Gap #14)
            await self._generate_causal_hypotheses(user_id, company_id, doc_id, facts)
        except Exception as e:
            logger.warning("Knowledge extraction failed: %s", e)

    async def _generate_causal_hypotheses(
        self,
        user_id: str,
        company_id: str,
        doc_id: str,
        facts: list[dict[str, Any]],
    ) -> None:
        """Generate causal hypotheses from document facts (Gap #14).

        Uses LLM to infer business implications from extracted facts,
        stored with confidence 0.50-0.60 as inferred knowledge.
        """
        if len(facts) < 2:
            return

        fact_text = "\n".join(f"- {f['fact']}" for f in facts[:10])
        prompt = (
            "Given these business facts from a document:\n\n"
            f"{fact_text}\n\n"
            "Generate 3-5 causal business hypotheses — inferences about "
            "what these facts imply for sales strategy. Format as JSON array:\n"
            "[\n"
            '  {"hypothesis": "Fact A → likely implication B", "confidence": 0.55}\n'
            "]\n"
            "Examples: 'Series C funding → hiring ramp → pipeline generation need', "
            "'FDA approval → commercial launch prep → need for KOL mapping'. "
            "Keep confidence between 0.50 and 0.60."
        )

        try:
            response = await self._llm.generate_response(
                messages=[{"role": "user", "content": prompt}],
                max_tokens=1000,
                temperature=0.5,
                task=TaskType.ENTITY_EXTRACT,
            )
            hypotheses: list[dict[str, Any]] = json.loads(response)
            for h in hypotheses:
                self._db.table("memory_semantic").insert(
                    {
                        "user_id": user_id,
                        "fact": h["hypothesis"],
                        "confidence": min(0.60, max(0.50, h.get("confidence", 0.55))),
                        "source": "inferred_during_document_upload",
                        "metadata": {
                            "category": "causal_hypothesis",
                            "document_id": doc_id,
                            "company_id": company_id,
                        },
                    }
                ).execute()
            logger.info(
                "Causal hypotheses generated from document",
                extra={
                    "document_id": doc_id,
                    "hypothesis_count": len(hypotheses),
                },
            )
        except Exception as e:
            logger.warning("Causal hypothesis generation failed: %s", e)

    async def _update_progress(self, doc_id: str, progress: float, status: str) -> None:
        """Update document processing progress.

        Args:
            doc_id: Document UUID.
            progress: Progress percentage (0-100).
            status: Processing status string.
        """
        self._db.table("company_documents").update(
            {
                "processing_progress": progress,
                "processing_status": status,
            }
        ).eq("id", doc_id).execute()

    async def _seed_entities_to_graph(
        self,
        entities: list[dict[str, str]],
        company_id: str,
        user_id: str,
        doc_id: str,
    ) -> None:
        """Seed extracted entities into Graphiti knowledge graph as episodes.

        Creates graph nodes/edges for entities found in documents so they're
        queryable via the knowledge graph, not just stored as JSON in chunks.

        Args:
            entities: List of entity dicts with 'name' and 'type' keys.
            company_id: Company UUID.
            user_id: User UUID.
            doc_id: Document UUID.
        """
        if not entities:
            return

        try:
            from src.memory.episodic import Episode, EpisodicMemory

            memory = EpisodicMemory()
            now = datetime.now(UTC)

            # Deduplicate by entity name
            seen: set[str] = set()
            unique_entities: list[dict[str, str]] = []
            for entity in entities:
                name = entity.get("name", "")
                if name and name not in seen:
                    seen.add(name)
                    unique_entities.append(entity)

            if not unique_entities:
                return

            # Build entity summary for graph seeding
            entity_summary = ", ".join(
                f"{e.get('name', '')} ({e.get('type', 'unknown')})" for e in unique_entities[:20]
            )

            episode = Episode(
                id=str(uuid.uuid4()),
                user_id=user_id,
                event_type="document_entity_extraction",
                content=(
                    f"Extracted {len(unique_entities)} entities from document: {entity_summary}"
                ),
                participants=[e.get("name", "") for e in unique_entities[:20]],
                occurred_at=now,
                recorded_at=now,
                context={
                    "document_id": doc_id,
                    "company_id": company_id,
                    "entity_count": len(unique_entities),
                    "entities": [
                        {"name": e.get("name", ""), "type": e.get("type", "")}
                        for e in unique_entities[:50]
                    ],
                },
            )
            await memory.store_episode(episode)

            logger.info(
                "Seeded %d entities to knowledge graph",
                len(unique_entities),
                extra={"document_id": doc_id, "entity_count": len(unique_entities)},
            )
        except Exception as e:
            logger.warning(
                "Failed to seed entities to knowledge graph: %s",
                e,
                extra={"document_id": doc_id},
            )

    async def _update_readiness(self, user_id: str, quality_score: float) -> None:
        """Update corporate_memory readiness based on document quality.

        Args:
            user_id: User UUID.
            quality_score: Document quality score (0-100).
        """
        try:
            from src.onboarding.orchestrator import OnboardingOrchestrator

            orch = OnboardingOrchestrator()
            boost = min(10.0, quality_score * 0.1)
            current = (
                self._db.table("onboarding_state")
                .select("readiness_scores")
                .eq("user_id", user_id)
                .maybe_single()
                .execute()
            )
            if current.data:
                scores = current.data.get("readiness_scores", {})
                new_score = min(100.0, scores.get("corporate_memory", 0) + boost)
                await orch.update_readiness_scores(user_id, {"corporate_memory": new_score})
        except Exception as e:
            logger.warning("Failed to update readiness: %s", e)

    async def _record_episodic(
        self,
        user_id: str,
        doc_id: str,
        chunk_count: int,
        entity_count: int,
        quality_score: float,
    ) -> None:
        """Record document processing event to episodic memory.

        Args:
            user_id: User UUID.
            doc_id: Document UUID.
            chunk_count: Number of chunks extracted.
            entity_count: Number of unique entities found.
            quality_score: Document quality score.
        """
        try:
            from src.memory.episodic import Episode, EpisodicMemory

            memory = EpisodicMemory()
            now = datetime.now(UTC)
            episode = Episode(
                id=str(uuid.uuid4()),
                user_id=user_id,
                event_type="onboarding_document_processed",
                content=str(
                    {
                        "document_id": doc_id,
                        "chunks": chunk_count,
                        "entities": entity_count,
                        "quality_score": quality_score,
                    }
                ),
                participants=[],
                occurred_at=now,
                recorded_at=now,
                context={
                    "document_id": doc_id,
                    "chunks": chunk_count,
                    "entities": entity_count,
                    "quality_score": quality_score,
                },
            )
            await memory.store_episode(episode)
        except Exception as e:
            logger.warning("Failed to record episodic event: %s", e)
